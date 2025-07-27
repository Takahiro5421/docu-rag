"""Utilities for extracting text from various file types."""

from pathlib import Path
from typing import Iterable

from io import BytesIO

try:
    import docx
    from openpyxl import load_workbook
    import PyPDF2
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependencies
    docx = None
    load_workbook = None
    PyPDF2 = None
    Presentation = None


def extract_text(file_path: Path) -> str:
    """Extract text from a single document.

    Supports .txt, .docx, .xlsx, .pdf and .pptx files.
    """

    if not file_path.exists():
        raise FileNotFoundError(file_path)

    suffix = file_path.suffix.lower()
    try:
        if suffix == ".txt":
            return file_path.read_text(encoding="utf-8")
        if suffix == ".docx" and docx:
            doc = docx.Document(str(file_path))
            return "\n".join(p.text for p in doc.paragraphs)
        if suffix == ".xlsx" and load_workbook:
            wb = load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    text_parts.append("\t".join(str(c) if c is not None else "" for c in row))
            return "\n".join(text_parts)
        if suffix == ".pdf" and PyPDF2:
            text_parts = []
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text_parts.append(page.extract_text() or "")
            return "\n".join(text_parts)
        if suffix == ".pptx" and Presentation:
            pres = Presentation(str(file_path))
            text_runs = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
    except Exception as e:  # pragma: no cover - runtime errors
        raise RuntimeError(f"Failed to extract text from {file_path}: {e}") from e

    raise ValueError(f"Unsupported file type: {file_path}")


def chunk_text(text: str, chunk_size: int = 500) -> Iterable[str]:
    """Yield fixed-size chunks from text."""
    for i in range(0, len(text), chunk_size):
        yield text[i : i + chunk_size]
